"""
03_bericht.py — Quartalsbericht (PPTX) + Dashboard-Excel + LV-Monitoring generieren.

Ausfuehrung: python 03_bericht.py
Voraussetzung: output/review/konsolidiert.xlsx + output/charts/ muessen existieren
Ergebnis:      output/reports/ (PPTX, Dashboard-Excel, LV-Monitoring)
"""

from pathlib import Path
from datetime import date
import copy

import pandas as pd
from openpyxl import Workbook
from openpyxl.styles import PatternFill, Font, Alignment, Border, Side, numbers
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from config import (
    PPTX_TEMPLATE, OUTPUT_REVIEW, OUTPUT_CHARTS, OUTPUT_REPORTS,
    QUARTAL, STICHTAG, AMPEL_FARBEN, DASHBOARD_COLS,
)


# =============================================================================
# 1. PPTX-Quartalsbericht
# =============================================================================

def _find_chart_pngs() -> dict[str, Path]:
    """Findet alle generierten Chart-PNGs."""
    charts = {}
    for f in OUTPUT_CHARTS.glob("*.png"):
        charts[f.stem] = f
    return charts


def _clear_slide_content(slide):
    """Entfernt alle Shapes ausser Titel von einer Folie."""
    shapes_to_remove = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip().lower()
            # Behalte Titel und Subtitel
            if shape.shape_id <= 2:
                continue
        shapes_to_remove.append(shape)
    for shape in shapes_to_remove:
        sp = shape._element
        sp.getparent().remove(sp)


def generate_pptx(df: pd.DataFrame):
    """Generiert den Quartalsbericht als PPTX basierend auf dem Template."""
    OUTPUT_REPORTS.mkdir(parents=True, exist_ok=True)

    # Template oeffnen
    if PPTX_TEMPLATE.exists():
        prs = Presentation(str(PPTX_TEMPLATE))
        print(f"  Template geladen: {PPTX_TEMPLATE.name}")
    else:
        prs = Presentation()
        print("  Kein Template gefunden, erstelle neue Praesentation.")

    charts = _find_chart_pngs()

    # --- Titelfolie aktualisieren (Folie 1) ---
    if len(prs.slides) > 0:
        slide = prs.slides[0]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if "Q4/2025" in run.text or "Q4 2025" in run.text:
                            run.text = run.text.replace("Q4/2025", QUARTAL).replace("Q4 2025", QUARTAL)
                        if "2025" in run.text and "Q" not in run.text and "LV" not in run.text:
                            run.text = run.text.replace("2025", "2026")

    # --- Chart-Folien: Bestehende Bilder ersetzen (Folien 5-10) ---
    chart_mapping = {
        4: "ampel_gesamt",           # Folie 5
        5: "ampel_nach_pag",         # Folie 6
        6: "projekte_nach_leistungsbereich",  # Folie 7
        7: "projekte_nach_phase",    # Folie 8
        8: "budget_plan_vs_ist",     # Folie 9
        9: "gantt_zeitleiste",       # Folie 10
    }

    for slide_idx, chart_name in chart_mapping.items():
        if slide_idx < len(prs.slides) and chart_name in charts:
            slide = prs.slides[slide_idx]
            # Bestehende Bilder entfernen
            pics_to_remove = []
            for shape in slide.shapes:
                if shape.shape_type == 13:  # Picture
                    pics_to_remove.append(shape)
            for pic in pics_to_remove:
                sp = pic._element
                sp.getparent().remove(sp)
            # Neues Chart-Bild einfuegen
            slide.shapes.add_picture(
                str(charts[chart_name]),
                Inches(0.5), Inches(1.5),
                Inches(9), Inches(5.5),
            )

    # --- Statusbericht-Folien: Gelbe/Rote Projekte (Folien 12-13) ---
    gelb_rot = df[df["ampelstatus"].isin(["Vorsicht", "Krise"])].copy()

    if len(prs.slides) > 11 and not gelb_rot.empty:
        # Folie 12: Tabelle mit gelben/roten Projekten
        slide = prs.slides[11]
        _clear_slide_content(slide)

        # Titel setzen
        if slide.shapes.title:
            slide.shapes.title.text = f"Statusberichte mit Handlungsbedarf ({QUARTAL})"

        # Tabelle erstellen
        rows_count = len(gelb_rot) + 1  # +1 fuer Header
        cols_count = 5
        table_shape = slide.shapes.add_table(
            rows_count, cols_count,
            Inches(0.3), Inches(1.5), Inches(9.4), Inches(0.4 * rows_count),
        )
        table = table_shape.table

        # Header
        headers = ["LV-Nr.", "Projektname", "Ampel", "Phase", "Risiko"]
        for j, h in enumerate(headers):
            cell = table.cell(0, j)
            cell.text = h
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(9)
                para.font.bold = True
                para.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x1F, 0x4E, 0x79)

        # Daten
        for i, (_, row) in enumerate(gelb_rot.iterrows(), 1):
            vals = [
                str(row.get("lv_nummer", "")),
                str(row.get("projektname", ""))[:50],
                str(row.get("ampelstatus", "")),
                str(row.get("projektphase", "")),
                str(row.get("risiko", "")),
            ]
            for j, val in enumerate(vals):
                cell = table.cell(i, j)
                cell.text = val
                for para in cell.text_frame.paragraphs:
                    para.font.size = Pt(8)
                # Ampelfarbe
                if j == 2:
                    ampel = row.get("ampelstatus", "")
                    if ampel == "Krise":
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(0xF4, 0x43, 0x36)
                    elif ampel == "Vorsicht":
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(0xFF, 0xC1, 0x07)

    # Speichern
    output_path = OUTPUT_REPORTS / f"Projektportfolio_LV-Vorhaben_{QUARTAL.replace('/', '')}.pptx"
    prs.save(str(output_path))
    print(f"  PPTX gespeichert: {output_path}")


# =============================================================================
# 2. Dashboard-Excel (Portfolio_Daten-Format)
# =============================================================================

def generate_dashboard_excel(df: pd.DataFrame):
    """Erzeugt ein Dashboard-Excel im Portfolio_Daten-Format mit Conditional Formatting."""
    OUTPUT_REPORTS.mkdir(parents=True, exist_ok=True)
    wb = Workbook()
    ws = wb.active
    ws.title = "Portfolio_Daten"

    # Spaltenreihenfolge gemaess DASHBOARD_COLS
    col_order = [name for _, name in sorted(DASHBOARD_COLS.items())]
    # Nur Spalten die im DataFrame existieren
    col_order = [c for c in col_order if c in df.columns]

    # Header (Zeile 1 = Titel, Daten ab Zeile 5 wie im Original)
    ws.merge_cells("A1:H1")
    title_cell = ws["A1"]
    title_cell.value = f"Projektportfolio LV-Vorhaben 2025-2027 — {QUARTAL}"
    title_cell.font = Font(size=14, bold=True, color="1F4E79")

    ws.merge_cells("A2:H2")
    ws["A2"].value = f"Stand: {STICHTAG.strftime('%d.%m.%Y')}"
    ws["A2"].font = Font(size=10, italic=True)

    # Spaltenheader in Zeile 4
    header_fill = PatternFill(start_color="1F4E79", end_color="1F4E79", fill_type="solid")
    header_font = Font(color="FFFFFF", bold=True, size=9)
    thin_border = Border(
        left=Side(style="thin"), right=Side(style="thin"),
        top=Side(style="thin"), bottom=Side(style="thin"),
    )

    # Lesbare Spaltennamen
    col_labels = {
        "lv_nummer": "LV-Nummer",
        "leistungsbereich": "Leistungsbereich",
        "kapitelzuordnung": "Kapitelzuordnung",
        "projektname": "Projektname/Bezeichnung",
        "kurzbeschreibung": "Kurzbeschreibung",
        "meilensteine": "Relevante Meilensteine",
        "projektleitung": "PL/Ansprechperson",
        "pag": "PAG",
        "organisationseinheiten": "Beteiligte OE",
        "lv_periode": "LV-Periode",
        "projektauftrag": "Projektauftrag",
        "psb_vorhanden": "PSB vorhanden",
        "lenkungsausschuss": "Lenkungsausschuss",
        "externe_kooperationspartner": "Externe Kooperation",
        "thema_begleitgespraech": "Thema Begleitgespraech",
        "budgeteinbehalt": "Budgeteinbehalt",
        "prioritaet": "Prioritaet",
        "risiko": "Risiko",
        "start": "Start",
        "ende": "Ende",
        "dauer_werktage": "Dauer (Werktage)",
        "kostenstelle": "Kostenstelle/IA",
        "plankosten_2025": "Plankosten 2025",
        "plankosten_2026": "Plankosten 2026",
        "plankosten_2027": "Plankosten 2027",
        "plankosten_gesamt": "Plankosten 25-27",
        "istkosten_2025": "Istkosten 2025",
        "istkosten_2026": "Istkosten 2026",
        "istkosten_2027": "Istkosten 2027",
        "istkosten_gesamt": "Ist-Kosten (Gesamt)",
        "istkosten_prozent": "Ist-Kosten in %",
        "projektphase": "Projektphase",
        "ampelstatus": "Ampelstatus Gesamt",
        "fertigstellungsgrad": "Fertigstellungsgrad (%)",
        "ziele": "Dazugehoerige Ziele",
        "status_aktuell": "Status (aktuell)",
        "risiken_aktuell": "Risiken (aktuell)",
        "entscheidung_aktuell": "Entscheidungsbedarf (aktuell)",
        "status_vorperiode": "Status (Vorperiode)",
        "risiken_vorperiode": "Risiken (Vorperiode)",
        "entscheidung_vorperiode": "Entscheidungsbedarf (Vorperiode)",
    }

    for col_idx, col_name in enumerate(col_order, 1):
        cell = ws.cell(row=4, column=col_idx)
        cell.value = col_labels.get(col_name, col_name)
        cell.fill = header_fill
        cell.font = header_font
        cell.alignment = Alignment(horizontal="center", wrap_text=True)
        cell.border = thin_border

    # Daten ab Zeile 5
    ampel_fills = {
        "In Ordnung": PatternFill(start_color="C6EFCE", end_color="C6EFCE", fill_type="solid"),
        "Vorsicht": PatternFill(start_color="FFEB9C", end_color="FFEB9C", fill_type="solid"),
        "Krise": PatternFill(start_color="FFC7CE", end_color="FFC7CE", fill_type="solid"),
    }

    ampel_col_idx = col_order.index("ampelstatus") + 1 if "ampelstatus" in col_order else None

    for row_idx, (_, row) in enumerate(df.iterrows(), 5):
        for col_idx, col_name in enumerate(col_order, 1):
            cell = ws.cell(row=row_idx, column=col_idx)
            val = row.get(col_name)
            if pd.isna(val):
                val = ""
            cell.value = val
            cell.border = thin_border
            cell.font = Font(size=9)

            # Prozentformat
            if col_name in ("istkosten_prozent", "fertigstellungsgrad"):
                if val and val != "":
                    cell.number_format = "0%"

            # Waehrungsformat
            if "kosten" in col_name and "prozent" not in col_name:
                cell.number_format = '#,##0.00 "EUR"'

        # Ampelstatus einfaerben
        if ampel_col_idx:
            ampel_cell = ws.cell(row=row_idx, column=ampel_col_idx)
            ampel_val = str(ampel_cell.value).strip()
            if ampel_val in ampel_fills:
                ampel_cell.fill = ampel_fills[ampel_val]

    # Spaltenbreiten
    widths = {"lv_nummer": 12, "projektname": 35, "kurzbeschreibung": 25, "pag": 10,
              "ampelstatus": 15, "projektphase": 20, "projektleitung": 18}
    for col_idx, col_name in enumerate(col_order, 1):
        ws.column_dimensions[get_column_letter(col_idx)].width = widths.get(col_name, 14)

    # Autofilter
    last_col = get_column_letter(len(col_order))
    last_row = 4 + len(df)
    ws.auto_filter.ref = f"A4:{last_col}{last_row}"

    # Zeile fixieren
    ws.freeze_panes = "A5"

    # Als Tabelle definieren (Tabelle1)
    from openpyxl.worksheet.table import Table, TableStyleInfo
    tab = Table(
        displayName="Tabelle1",
        ref=f"A4:{last_col}{last_row}",
    )
    tab.tableStyleInfo = TableStyleInfo(
        name="TableStyleMedium2", showFirstColumn=False,
        showLastColumn=False, showRowStripes=True, showColumnStripes=False,
    )
    ws.add_table(tab)

    output_path = OUTPUT_REPORTS / f"Portfolio_Dashboard_{QUARTAL.replace('/', '')}.xlsx"
    wb.save(str(output_path))
    print(f"  Dashboard-Excel gespeichert: {output_path}")


# =============================================================================
# 3. LV-Monitoring
# =============================================================================

def generate_lv_monitoring(df: pd.DataFrame):
    """Erzeugt das LV-Monitoring-Excel: Ampelstatus + Erlaeuterung + Zielwerte pro Vorhaben."""
    OUTPUT_REPORTS.mkdir(parents=True, exist_ok=True)
    wb = Workbook()
    ws = wb.active
    ws.title = "LV-Monitoring"

    # Titel
    ws.merge_cells("A1:H1")
    ws["A1"].value = f"LV-Monitoring — {QUARTAL}"
    ws["A1"].font = Font(size=14, bold=True, color="1F4E79")

    ws.merge_cells("A2:H2")
    ws["A2"].value = f"Stand: {STICHTAG.strftime('%d.%m.%Y')} | Leistungsvereinbarung 2025-2027"
    ws["A2"].font = Font(size=10, italic=True)

    # Header
    headers = [
        "LV-Nummer", "Projektname", "Leistungsbereich", "PAG",
        "Ampelstatus", "Erlaeuterung Ampel",
        "Ziele gemaess LV", "Fertigstellungsgrad",
    ]
    header_fill = PatternFill(start_color="1F4E79", end_color="1F4E79", fill_type="solid")
    header_font = Font(color="FFFFFF", bold=True, size=10)
    thin_border = Border(
        left=Side(style="thin"), right=Side(style="thin"),
        top=Side(style="thin"), bottom=Side(style="thin"),
    )

    for col_idx, h in enumerate(headers, 1):
        cell = ws.cell(row=4, column=col_idx)
        cell.value = h
        cell.fill = header_fill
        cell.font = header_font
        cell.alignment = Alignment(horizontal="center", wrap_text=True)
        cell.border = thin_border

    # Daten
    ampel_fills = {
        "In Ordnung": PatternFill(start_color="C6EFCE", end_color="C6EFCE", fill_type="solid"),
        "Vorsicht": PatternFill(start_color="FFEB9C", end_color="FFEB9C", fill_type="solid"),
        "Krise": PatternFill(start_color="FFC7CE", end_color="FFC7CE", fill_type="solid"),
    }

    for row_idx, (_, row) in enumerate(df.iterrows(), 5):
        vals = [
            row.get("lv_nummer", ""),
            row.get("projektname", ""),
            row.get("leistungsbereich", ""),
            row.get("pag", ""),
            row.get("ampelstatus", ""),
            row.get("status_aktuell", ""),
            row.get("ziele", ""),
            row.get("fertigstellungsgrad", 0),
        ]
        for col_idx, val in enumerate(vals, 1):
            cell = ws.cell(row=row_idx, column=col_idx)
            if pd.isna(val):
                val = ""
            cell.value = val
            cell.border = thin_border
            cell.font = Font(size=9)
            cell.alignment = Alignment(wrap_text=True, vertical="top")

            # Ampelstatus einfaerben
            if col_idx == 5:
                ampel_val = str(val).strip()
                if ampel_val in ampel_fills:
                    cell.fill = ampel_fills[ampel_val]

            # Fertigstellungsgrad als Prozent
            if col_idx == 8:
                cell.number_format = "0%"

    # Spaltenbreiten
    for col_idx, width in enumerate([12, 35, 20, 10, 15, 45, 35, 12], 1):
        ws.column_dimensions[get_column_letter(col_idx)].width = width

    ws.freeze_panes = "A5"

    output_path = OUTPUT_REPORTS / f"LV_Monitoring_{QUARTAL.replace('/', '')}.xlsx"
    wb.save(str(output_path))
    print(f"  LV-Monitoring gespeichert: {output_path}")


# =============================================================================
# Main
# =============================================================================

def main():
    print(f"=== Berichtsgenerierung ({QUARTAL}) ===\n")

    # Daten laden
    data_path = OUTPUT_REVIEW / "konsolidiert.xlsx"
    if not data_path.exists():
        print("FEHLER: Bitte zuerst 01_konsolidierung.py ausfuehren!")
        return
    df = pd.read_excel(data_path, sheet_name="Konsolidiert")
    print(f"  {len(df)} Projekte geladen.\n")

    # 1. PPTX
    print("1. Quartalsbericht (PPTX)...")
    generate_pptx(df)
    print()

    # 2. Dashboard-Excel
    print("2. Dashboard-Excel...")
    generate_dashboard_excel(df)
    print()

    # 3. LV-Monitoring
    print("3. LV-Monitoring...")
    generate_lv_monitoring(df)

    print(f"\nFertig. Berichte in: {OUTPUT_REPORTS}")


if __name__ == "__main__":
    main()
